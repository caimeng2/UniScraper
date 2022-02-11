"""A scraper that extracts text from multiple types of webpages, including html, pdf, word documents, presentation slides, and spreadsheets"""
import io
import csv
import time
import pandas as pd
import requests
from bs4 import BeautifulSoup
from bs4.element import Comment
import nltk
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from webdriver_manager.chrome import ChromeDriverManager
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.layout import LAParams
from pdfminer.converter import TextConverter
from pdfminer.pdfpage import PDFPage
import docx
from pptx import Presentation

options = Options()
options.add_argument('--headless')
driver = webdriver.Chrome(ChromeDriverManager().install(), options=options)

def tag_visible(element):
    """A helper function to filter visible tags"""
    if element.parent.name in ['style', 'script', 'head', 'title', 'meta', '[document]']:
        return False
    if isinstance(element, Comment):
        return False
    return True

def pdf_to_text(pdf_file):
    """Extract text from pdf file"""
    text_memory_file = io.StringIO()
    rsrcmgr = PDFResourceManager()
    device = TextConverter(rsrcmgr, text_memory_file, laparams=LAParams())
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    for page in PDFPage.get_pages(pdf_file):
        interpreter.process_page(page)
    text = text_memory_file.getvalue()
    text_memory_file.close()
    return text

def doc_to_text(doc_file):
    """Extract text from word document"""
    doc = docx.Document(doc_file)
    text = []
    for i in range(len(doc.paragraphs)):
        para = doc.paragraphs[i]
        text.append(para.text)
    return "\n\n".join(text)

def ppt_to_text(ppt_file):
    """Extract text from presentation slides"""
    ppt = Presentation(ppt_file)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)

def excel_to_text(excel_file):
    """Extract text from spreadsheets"""
    data = pd.ExcelFile(excel_file)
    text = []
    for i in range(len(data.sheet_names)):
        sheet = data.parse(data.sheet_names[i])
        for content in sheet.values.tolist():
            text.append(content[0])
    return "\n\n".join(text)

def text_from_url(url):
    """Extract text from web pages"""
    try:
        req = requests.get(url)
        if 'text/html' in req.headers["Content-Type"]:
            driver.get(url)
            time.sleep(15)
            body = driver.page_source
            soup = BeautifulSoup(body, 'html.parser')
            texts = soup.findAll(text=True)
            visible_texts = filter(tag_visible, texts)
            text = u" ".join(t.strip() for t in visible_texts)
        elif 'application/pdf' in req.headers["Content-Type"]:
            pdf_memory_file = io.BytesIO()
            pdf_memory_file.write(req.content)
            text = pdf_to_text(pdf_memory_file)
        elif 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in req.headers["Content-Type"]:
            word_memory_file = io.BytesIO()
            word_memory_file.write(req.content)
            text = doc_to_text(word_memory_file)
        elif 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in req.headers["Content-Type"]:
            ppt_memory_file = io.BytesIO()
            ppt_memory_file.write(req.content)
            text = ppt_to_text(ppt_memory_file)
        elif 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in req.headers["Content-Type"]:
            excel_memory_file = io.BytesIO()
            excel_memory_file.write(req.content)
            text = excel_to_text(excel_memory_file)
        else:
            csv.writer(open("error.csv", "a")).writerow(["Unsupported content type: ", req.headers["Content-Type"], url])
            text = ""
    except Exception as err:
        csv.writer(open("error.csv", "a")).writerow([err, url])
        text = ""
    return text

def remove_non_eng(text):
    """A helper funtion to remove non-english text"""
    words = set(nltk.corpus.words.words())
    english = " ".join(w for w in nltk.wordpunct_tokenize(text) if w.lower() in words or not w.isalpha())
    return english

def paragraph_from_text(text, search_string):
    """Get paragraphs containing a keyword with its context"""
    split_text = text.lower().split("\n\n")
    found = False
    previous = ""
    para = []
    # to extract paragraphs containing a keyword and its immediate previous and after paragraphs
    for p in split_text:
        if found:
            para.append(p)
            found = False
        if p.find(search_string.lower()) != -1:
            para.append(previous)
            para.append(p)
            found = True
        else:
            found = False
        previous = p
    para = list(set(para)) # remove duplicates
    para_str = " ".join(para)
    return para_str

class uniscraper():
    """A scraper that works with multiple types of webpages"""
    def __init__(self, url):
        self.url = url
        self.text = text_from_url(url)
        self.english = remove_non_eng(self.text)
        self.para = None
    def search(self, string):
        """Search text for a keyword"""
        self.para = paragraph_from_text(self.text, string)
        return self.para

    def highlight_search(self, string):
        """Search text for a keyword"""
        print((f"\033[91m" + string + f"\033[0m").join(self.text.split(string)))
